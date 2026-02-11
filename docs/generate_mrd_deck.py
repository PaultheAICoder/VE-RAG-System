#!/usr/bin/env python3
"""Generate MRD presentation deck for CEO meeting."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x00, 0xB4, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MED_GRAY = RGBColor(0x99, 0x99, 0xAA)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x3D)
RED_ACCENT = RGBColor(0xE8, 0x4D, 0x4D)
GOLD_ACCENT = RGBColor(0xFF, 0xB7, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_dark_bg(slide):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_light_bg(slide):
    """Add light background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_accent_bar(slide, top=Inches(0), color=ACCENT_BLUE, height=Inches(0.06)):
    """Add thin accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box with single text run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tf


def add_bullet_list(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=16,
    color=DARK_TEXT,
    bold_prefix=True,
    spacing=Pt(8),
):
    """Add a bulleted list. Items can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        if isinstance(item, tuple):
            run_bold = p.add_run()
            run_bold.text = item[0]
            run_bold.font.size = Pt(font_size)
            run_bold.font.bold = True
            run_bold.font.color.rgb = color
            run_bold.font.name = "Calibri"
            run_rest = p.add_run()
            run_rest.text = item[1]
            run_rest.font.size = Pt(font_size)
            run_rest.font.bold = False
            run_rest.font.color.rgb = color
            run_rest.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = f"  {item}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return tf


def add_card(slide, left, top, width, height, title, body, title_color=ACCENT_BLUE):
    """Add a rounded rectangle card with title and body."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()

    # Title
    add_text_box(
        slide,
        left + Inches(0.2),
        top + Inches(0.15),
        width - Inches(0.4),
        Inches(0.4),
        title,
        font_size=15,
        bold=True,
        color=title_color,
    )

    # Body
    add_text_box(
        slide,
        left + Inches(0.2),
        top + Inches(0.55),
        width - Inches(0.4),
        height - Inches(0.7),
        body,
        font_size=13,
        color=DARK_TEXT,
    )


# =============================================================================
# SLIDE 1: Title
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)
add_accent_bar(slide, top=Inches(3.2), color=ACCENT_BLUE, height=Inches(0.04))

add_text_box(
    slide,
    Inches(1),
    Inches(1.5),
    Inches(11),
    Inches(1.2),
    "AI Ready RAG",
    font_size=54,
    bold=True,
    color=WHITE,
)
add_text_box(
    slide,
    Inches(1),
    Inches(2.5),
    Inches(11),
    Inches(0.6),
    "Marketing Requirements Document",
    font_size=28,
    color=ACCENT_BLUE,
)

add_text_box(
    slide,
    Inches(1),
    Inches(3.8),
    Inches(11),
    Inches(1.5),
    "Air-Gapped Enterprise Knowledge Assistant\nfor NVIDIA DGX Spark & Hosted Deployment",
    font_size=20,
    color=MED_GRAY,
)

add_text_box(
    slide,
    Inches(1),
    Inches(5.8),
    Inches(5),
    Inches(0.8),
    "February 13, 2026  |  Vital Enterprises — Internal",
    font_size=14,
    color=MED_GRAY,
)

add_text_box(
    slide,
    Inches(7),
    Inches(5.8),
    Inches(5),
    Inches(0.8),
    "Presented by Jason & Paul",
    font_size=14,
    color=MED_GRAY,
    alignment=PP_ALIGN.RIGHT,
)

# =============================================================================
# SLIDE 2: Executive Summary
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Executive Summary",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.1),
    Inches(11.5),
    Inches(0.8),
    "AI Ready RAG is an air-gapped enterprise knowledge assistant that enables organizations to search, "
    "query, and get cited answers from internal documents — without any data leaving their infrastructure.",
    font_size=18,
    color=DARK_TEXT,
)

# Three differentiators
card_y = Inches(2.2)
card_h = Inches(1.8)
card_w = Inches(3.6)
gap = Inches(0.3)

add_card(
    slide,
    Inches(0.8),
    card_y,
    card_w,
    card_h,
    "Air-Gapped Security",
    "No data leaves the customer's network.\nFull runtime independence\nfrom the internet.",
)

add_card(
    slide,
    Inches(0.8) + card_w + gap,
    card_y,
    card_w,
    card_h,
    "Lower Total Cost",
    "Predictable $850/month pricing vs.\nenterprise platforms at\n$2,000-$8,600+/month.",
)

add_card(
    slide,
    Inches(0.8) + 2 * (card_w + gap),
    card_y,
    card_w,
    card_h,
    "Customer Service",
    "Hands-on rollout, operator\nonboarding, and ongoing support.\nNot self-service DIY.",
)

# Timeline
add_text_box(
    slide,
    Inches(0.8),
    Inches(4.3),
    Inches(11),
    Inches(0.5),
    "Business Targets",
    font_size=20,
    bold=True,
    color=DARK_TEXT,
)

milestones = [
    ("Mar 16 ", " — Internal testing (2-3 Vital Enterprises organizations)"),
    ("Apr 20 ", " — Demo customers enabled"),
    ("May 25 ", " — Product launch (GA)"),
    ("Sep 16 ", " — 6 paid pilots closed ($5,100 MRR)"),
]
add_bullet_list(
    slide,
    Inches(0.8),
    Inches(4.9),
    Inches(11),
    Inches(2.2),
    milestones,
    font_size=17,
    color=DARK_TEXT,
)

# =============================================================================
# SLIDE 3: The Problem
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "The Problem",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.0),
    Inches(11),
    Inches(0.6),
    "Critical operational knowledge is trapped in documents employees can't find quickly.",
    font_size=20,
    color=DARK_TEXT,
)

problems = [
    (
        "Document search is slow",
        " — Employees spend 20-30% of time searching across scattered systems",
    ),
    (
        "Onboarding takes too long",
        " — New hires take weeks; institutional knowledge lives in people's heads",
    ),
    (
        "Process execution is inconsistent",
        " — Different employees follow different procedure versions",
    ),
    ("Knowledge lost on turnover", " — When experienced staff leave, their knowledge leaves too"),
    (
        "Cloud AI is not acceptable",
        " — Security/compliance prevents sending documents to cloud services",
    ),
]

y_pos = Inches(1.8)
for title, desc in problems:
    card_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.5), Inches(0.85)
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = LIGHT_GRAY
    card_shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1.1), y_pos + Inches(0.1), Inches(11), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_b = p.add_run()
    run_b.text = title
    run_b.font.size = Pt(16)
    run_b.font.bold = True
    run_b.font.color.rgb = ACCENT_BLUE
    run_b.font.name = "Calibri"
    run_r = p.add_run()
    run_r.text = desc
    run_r.font.size = Pt(15)
    run_r.font.color.rgb = DARK_TEXT
    run_r.font.name = "Calibri"

    y_pos += Inches(0.95)

# =============================================================================
# SLIDE 4: Why Now
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Why Now",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.0),
    Inches(11),
    Inches(0.6),
    "Four market inflections create a unique window (March - September 2026):",
    font_size=20,
    color=DARK_TEXT,
)

why_now = [
    ("Hardware Cost Inflection", "DGX Spark at $3,999 makes local\nGPU inference viable for SMB"),
    (
        "Model Quality Inflection",
        "8B open-source models now rival\ncloud models from 12 months ago",
    ),
    (
        "Regulatory Pressure",
        "Increasing data sovereignty requirements\nin legal, financial, government",
    ),
    (
        "Competitive Gap",
        "Competitors are cloud-first or DIY.\nTrue air-gap + turnkey is unoccupied.",
    ),
]

card_w = Inches(2.7)
gap = Inches(0.25)
start_x = Inches(0.8)
for i, (title, body) in enumerate(why_now):
    x = start_x + i * (card_w + gap)
    add_card(slide, x, Inches(1.9), card_w, Inches(2.2), title, body, title_color=ACCENT_GREEN)

# =============================================================================
# SLIDE 5: Product Vision + Dual Deployment
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Product Vision & Deployment",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.0),
    Inches(11),
    Inches(0.8),
    "Every mid-market and regulated enterprise has critical operational knowledge trapped in documents.\n"
    "AI Ready RAG makes it instantly accessible — all running inside the customer's own walls.",
    font_size=18,
    color=DARK_TEXT,
)

# Two deployment cards
dep_w = Inches(5.5)
dep_h = Inches(3.4)

# On-Premise card
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), dep_w, dep_h)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(1.1),
    Inches(2.4),
    Inches(5),
    Inches(0.4),
    "On-Premise (DGX Spark)",
    font_size=18,
    bold=True,
    color=ACCENT_BLUE,
)

on_prem_items = [
    "Customer purchases NVIDIA DGX Spark ($3,999)",
    "AI Ready RAG pre-installed or deployed on-site",
    "All data stays on customer hardware",
    "Zero network dependency at runtime",
    "Physical air-gap for maximum security",
]
add_bullet_list(
    slide,
    Inches(1.1),
    Inches(2.9),
    Inches(5),
    Inches(2.5),
    on_prem_items,
    font_size=14,
    color=DARK_TEXT,
)

# Hosted card
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), dep_w, dep_h)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(7.1),
    Inches(2.4),
    Inches(5),
    Inches(0.4),
    "Hosted (Dedicated Instance)",
    font_size=18,
    bold=True,
    color=ACCENT_BLUE,
)

hosted_items = [
    "Vital Enterprises manages the infrastructure",
    "Dedicated single-tenant instance per customer",
    "NOT multi-tenant — fully isolated",
    "Private encrypted connection only",
    "Customer data never exposed to public internet",
]
add_bullet_list(
    slide,
    Inches(7.1),
    Inches(2.9),
    Inches(5),
    Inches(2.5),
    hosted_items,
    font_size=14,
    color=DARK_TEXT,
)

# Bottom callout
add_text_box(
    slide,
    Inches(0.8),
    Inches(5.9),
    Inches(11.5),
    Inches(0.8),
    "Both models deliver the same data isolation guarantee. "
    "The difference is who manages the hardware, not how the data is protected.",
    font_size=16,
    bold=True,
    color=ACCENT_BLUE,
    alignment=PP_ALIGN.CENTER,
)

# =============================================================================
# SLIDE 6: Target Customer
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Target Customer",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

# ICP card
icp_items = [
    ("Company size: ", "20-300 employees"),
    ("Verticals: ", "Legal, Manufacturing, Financial Services, Agriculture, Landscaping"),
    ("Geography: ", "West Coast first (Portland, OR hub)"),
    ("Workflow: ", "Document-heavy operations (SOPs, policies, training, compliance)"),
    ("Security: ", "Cloud hesitation or air-gap requirement"),
    ("Buyer: ", "Operations leader (not CTO or ML team)"),
    ("Budget: ", "Can approve $850/month without enterprise procurement cycle"),
]
add_bullet_list(
    slide,
    Inches(0.8),
    Inches(1.1),
    Inches(6),
    Inches(4.5),
    icp_items,
    font_size=16,
    color=DARK_TEXT,
)

# Primary persona card
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.1), Inches(5.2), Inches(4.8)
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(7.6),
    Inches(1.3),
    Inches(4.6),
    Inches(0.4),
    "Primary Buyer: Operations Leader",
    font_size=16,
    bold=True,
    color=ACCENT_BLUE,
)

persona_items = [
    "VP Ops / Director of Operations",
    "10-100 direct/indirect reports",
    "Pain: search time, slow onboarding",
    "Pain: inconsistent processes",
    "Pain: knowledge loss from turnover",
    "Controls departmental OpEx",
    "Not technical enough for DIY RAG",
]
add_bullet_list(
    slide,
    Inches(7.6),
    Inches(1.8),
    Inches(4.6),
    Inches(3.8),
    persona_items,
    font_size=14,
    color=DARK_TEXT,
)

# =============================================================================
# SLIDE 7: MVP Features
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "MVP Features",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

# Implemented column
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(3.7), Inches(5.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.fill.background()

add_text_box(
    slide,
    Inches(1.0),
    Inches(1.3),
    Inches(3.3),
    Inches(0.4),
    "IMPLEMENTED",
    font_size=18,
    bold=True,
    color=ACCENT_GREEN,
)

impl_items = [
    "Document ingestion (6 formats)",
    "Natural language query + citations",
    "Pre-retrieval access control",
    "Multi-turn chat sessions",
    "User management + tags",
    "React web application",
    "3-level audit logging",
    "JWT authentication",
]
add_bullet_list(
    slide,
    Inches(1.0),
    Inches(1.8),
    Inches(3.3),
    Inches(4.5),
    impl_items,
    font_size=14,
    color=DARK_TEXT,
)

# Planned column
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.2), Inches(3.7), Inches(5.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
shape.line.fill.background()

add_text_box(
    slide,
    Inches(5.0),
    Inches(1.3),
    Inches(3.3),
    Inches(0.4),
    "PLANNED (P0)",
    font_size=18,
    bold=True,
    color=GOLD_ACCENT,
)

planned_items = [
    "AD/SSO integration (Entra ID)",
    "Embeddable chat widget",
    "Usage reports & analytics",
    "Backup / restore",
    "Tiered XLSX ingestion",
]
add_bullet_list(
    slide,
    Inches(5.0),
    Inches(1.8),
    Inches(3.3),
    Inches(4.5),
    planned_items,
    font_size=14,
    color=DARK_TEXT,
)

# In Development column
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.2), Inches(3.7), Inches(5.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
shape.line.fill.background()

add_text_box(
    slide,
    Inches(9.0),
    Inches(1.3),
    Inches(3.3),
    Inches(0.4),
    "IN DEVELOPMENT",
    font_size=18,
    bold=True,
    color=ACCENT_BLUE,
)

dev_items = [
    "RAG quality evaluation (RAGAS)",
    "Confidence scoring (CITE/ROUTE)",
    "Hallucination checking",
    "Tiered PDF ingestion",
]
add_bullet_list(
    slide,
    Inches(9.0),
    Inches(1.8),
    Inches(3.3),
    Inches(4.5),
    dev_items,
    font_size=14,
    color=DARK_TEXT,
)

# Doc formats footer
add_text_box(
    slide,
    Inches(0.8),
    Inches(6.85),
    Inches(11.5),
    Inches(0.4),
    "Supported formats: PDF, DOCX, XLSX, PPTX, TXT, MD  |  English (MVP)  |  Spanish (evaluating)",
    font_size=13,
    color=MED_GRAY,
    alignment=PP_ALIGN.CENTER,
)

# =============================================================================
# SLIDE 8: Competitive Quadrant
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Competitive Landscape",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.0),
    Inches(11),
    Inches(0.5),
    "15 competitors analyzed. AI Ready RAG occupies an uncontested position: affordable + air-gapped + turnkey.",
    font_size=18,
    color=DARK_TEXT,
)

# Quadrant background
quad_left = Inches(1.5)
quad_top = Inches(1.8)
quad_w = Inches(10)
quad_h = Inches(5.0)

# Horizontal axis
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, quad_left, quad_top + quad_h // 2, quad_w, Inches(0.02)
)
shape.fill.solid()
shape.fill.fore_color.rgb = MED_GRAY
shape.line.fill.background()

# Vertical axis
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, quad_left + quad_w // 2, quad_top, Inches(0.02), quad_h
)
shape.fill.solid()
shape.fill.fore_color.rgb = MED_GRAY
shape.line.fill.background()

# Axis labels
add_text_box(
    slide,
    quad_left + quad_w // 2 - Inches(1.5),
    quad_top - Inches(0.35),
    Inches(3),
    Inches(0.3),
    "High Air-Gap Capability",
    font_size=12,
    bold=True,
    color=MED_GRAY,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    quad_left + quad_w // 2 - Inches(1.2),
    quad_top + quad_h + Inches(0.05),
    Inches(3),
    Inches(0.3),
    "Low Air-Gap Capability",
    font_size=12,
    bold=True,
    color=MED_GRAY,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    quad_left - Inches(0.3),
    quad_top + quad_h // 2 - Inches(0.35),
    Inches(2),
    Inches(0.3),
    "SMB Affordable",
    font_size=12,
    bold=True,
    color=MED_GRAY,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    quad_left + quad_w - Inches(1.8),
    quad_top + quad_h // 2 - Inches(0.35),
    Inches(2),
    Inches(0.3),
    "Enterprise Priced",
    font_size=12,
    bold=True,
    color=MED_GRAY,
    alignment=PP_ALIGN.CENTER,
)

# AI Ready RAG - top left (affordable + air-gapped) with star
star = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    quad_left + Inches(0.5),
    quad_top + Inches(0.3),
    Inches(2.5),
    Inches(0.7),
)
star.fill.solid()
star.fill.fore_color.rgb = ACCENT_BLUE
star.line.fill.background()
add_text_box(
    slide,
    quad_left + Inches(0.6),
    quad_top + Inches(0.35),
    Inches(2.3),
    Inches(0.5),
    "AI Ready RAG",
    font_size=16,
    bold=True,
    color=WHITE,
    alignment=PP_ALIGN.CENTER,
)

# OSS Stacks - bottom left
add_text_box(
    slide,
    quad_left + Inches(0.3),
    quad_top + quad_h // 2 + Inches(0.5),
    Inches(3.5),
    Inches(0.8),
    "OSS Stacks\n(Qdrant, OpenSearch, Haystack,\nLangChain, LlamaIndex)",
    font_size=13,
    color=DARK_TEXT,
)

# Enterprise - top right
add_text_box(
    slide,
    quad_left + quad_w // 2 + Inches(1),
    quad_top + Inches(0.3),
    Inches(3.5),
    Inches(0.6),
    "Sinequa, Mindbreeze\n(enterprise price)",
    font_size=13,
    color=DARK_TEXT,
)

# Cloud / Mid - bottom right
add_text_box(
    slide,
    quad_left + quad_w // 2 + Inches(0.8),
    quad_top + quad_h // 2 + Inches(0.5),
    Inches(4),
    Inches(0.8),
    "AWS Kendra, Azure AI Search,\nElastic, IBM watsonx, SearchBlox",
    font_size=13,
    color=DARK_TEXT,
)

# =============================================================================
# SLIDE 9: Competitive Pricing
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Competitive Pricing",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.0),
    Inches(11),
    Inches(0.5),
    "AI Ready RAG at $850/month is the only affordable, turnkey, air-gapped option.",
    font_size=20,
    color=DARK_TEXT,
)

# Price comparison as horizontal bars
competitors = [
    ("AI Ready RAG", 850, True),
    ("AWS Kendra", 230, False),
    ("Weaviate", 400, False),
    ("SearchBlox", 2083, False),
    ("IBM watsonx", 1050, False),
    ("Mindbreeze", 8641, False),
]

max_price = 8641
bar_area_left = Inches(3.2)
bar_area_width = Inches(8.5)
bar_height = Inches(0.55)
start_y = Inches(1.8)

for i, (name, price, is_us) in enumerate(competitors):
    y = start_y + i * Inches(0.85)

    # Label
    add_text_box(
        slide,
        Inches(0.8),
        y,
        Inches(2.2),
        bar_height,
        name,
        font_size=15,
        bold=is_us,
        color=ACCENT_BLUE if is_us else DARK_TEXT,
        alignment=PP_ALIGN.RIGHT,
    )

    # Bar
    bar_width_pct = min(price / max_price, 1.0)
    bar_w = int(bar_area_width * bar_width_pct)
    if bar_w < Inches(0.5):
        bar_w = Inches(0.5)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        bar_area_left,
        y + Inches(0.05),
        bar_w,
        bar_height - Inches(0.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE if is_us else RGBColor(0xCC, 0xCC, 0xDD)
    bar.line.fill.background()

    # Price label
    price_str = f"${price:,}/mo" if price < 10000 else f"${price:,}/mo"
    label_x = bar_area_left + bar_w + Inches(0.15)
    add_text_box(
        slide,
        label_x,
        y,
        Inches(1.5),
        bar_height,
        price_str,
        font_size=14,
        bold=is_us,
        color=ACCENT_BLUE if is_us else DARK_TEXT,
    )

# Notes
notes = [
    "OSS options (Qdrant, OpenSearch, LangChain) are $0 license but require months of engineering to productionize.",
    "Cloud options (AWS, Azure, GCP) lack true air-gap capability.",
]
add_bullet_list(
    slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1.2), notes, font_size=13, color=MED_GRAY
)

# =============================================================================
# SLIDE 10: Pricing & Packaging
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Pricing & Packaging",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

# Pilot package card
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5.5), Inches(4.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(1.1),
    Inches(1.4),
    Inches(5),
    Inches(0.4),
    "Pilot Package",
    font_size=22,
    bold=True,
    color=ACCENT_BLUE,
)

add_text_box(
    slide,
    Inches(1.1),
    Inches(1.9),
    Inches(5),
    Inches(0.6),
    "$850/month  |  6 months  |  $5,100 total",
    font_size=20,
    bold=True,
    color=DARK_TEXT,
)

pilot_includes = [
    "Setup and integration support",
    "Initial knowledge base ingestion",
    "Operator onboarding",
    "Monthly optimization review",
    "35% minimum margin target",
]
add_bullet_list(
    slide,
    Inches(1.1),
    Inches(2.6),
    Inches(5),
    Inches(3),
    pilot_includes,
    font_size=15,
    color=DARK_TEXT,
)

# Hardware card
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.5), Inches(2.0)
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(7.1),
    Inches(1.4),
    Inches(5),
    Inches(0.4),
    "Hardware (On-Prem Only)",
    font_size=18,
    bold=True,
    color=ACCENT_BLUE,
)

add_text_box(
    slide,
    Inches(7.1),
    Inches(1.9),
    Inches(5),
    Inches(0.4),
    "NVIDIA DGX Spark: $3,999 (customer purchase)",
    font_size=15,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(7.1),
    Inches(2.4),
    Inches(5),
    Inches(0.5),
    "Bundling model TBD — options: separate purchase,\nbundled pricing, or lease-to-own",
    font_size=13,
    color=RED_ACCENT,
)

# Decisions needed card
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.5), Inches(5.5), Inches(2.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
shape.line.fill.background()

add_text_box(
    slide,
    Inches(7.1),
    Inches(3.7),
    Inches(5),
    Inches(0.4),
    "DECISIONS NEEDED",
    font_size=16,
    bold=True,
    color=RED_ACCENT,
)

decision_items = [
    "Hardware bundling model",
    "Post-pilot conversion pricing",
    "Tiered support levels",
    "Hosted instance cost model",
]
add_bullet_list(
    slide,
    Inches(7.1),
    Inches(4.2),
    Inches(5),
    Inches(1.3),
    decision_items,
    font_size=14,
    color=DARK_TEXT,
)

# =============================================================================
# SLIDE 11: Go-to-Market
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Go-to-Market Strategy",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

# Left: Messaging
add_text_box(
    slide,
    Inches(0.8),
    Inches(1.1),
    Inches(5.5),
    Inches(0.4),
    "Messaging & Channels",
    font_size=20,
    bold=True,
    color=ACCENT_BLUE,
)

gtm_items = [
    (
        "Primary message: ",
        '"Air-gapped knowledge assistant with predictable SMB pricing and hands-on support."',
    ),
    ("Launch market: ", "West Coast, Portland OR hub"),
    ("Primary channel: ", "Direct outbound (email + LinkedIn + warm intros)"),
    ("Secondary: ", "Local events, MSP/IT consultant referrals"),
    ("Budget: ", "<$10,000 total through Sep 16"),
    ("Max CAC: ", "$1,500 per paid pilot"),
]
add_bullet_list(
    slide,
    Inches(0.8),
    Inches(1.6),
    Inches(5.8),
    Inches(3.5),
    gtm_items,
    font_size=15,
    color=DARK_TEXT,
)

# Right: Verticals
add_text_box(
    slide,
    Inches(7.3),
    Inches(1.1),
    Inches(5),
    Inches(0.4),
    "Priority Verticals",
    font_size=20,
    bold=True,
    color=ACCENT_BLUE,
)

verticals = [
    ("1. Legal", " — document-heavy, strict confidentiality"),
    ("2. Manufacturing", " — SOPs, safety, compliance docs"),
    ("3. Financial Services", " — regulatory, audit, data sovereignty"),
    ("4. Agriculture", " — process docs, seasonal onboarding"),
    ("5. Landscaping", " — operational SOPs, turnover mgmt"),
]
add_bullet_list(
    slide,
    Inches(7.3),
    Inches(1.6),
    Inches(5),
    Inches(3.5),
    verticals,
    font_size=15,
    color=DARK_TEXT,
)

# Win themes at bottom
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(1.1),
    Inches(5.4),
    Inches(5),
    Inches(0.3),
    "Win Themes",
    font_size=16,
    bold=True,
    color=ACCENT_GREEN,
)

win_items = [
    ("Security: ", '"Built for environments where internet-independent runtime matters."'),
    ("Cost: ", '"Predictable pricing and lower total ownership burden for SMB operations."'),
    ("Service: ", '"Hands-on rollout and ongoing support instead of self-service complexity."'),
]
add_bullet_list(
    slide,
    Inches(1.1),
    Inches(5.7),
    Inches(11),
    Inches(1.0),
    win_items,
    font_size=13,
    color=DARK_TEXT,
    spacing=Pt(4),
)

# =============================================================================
# SLIDE 12: Timeline
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Timeline & Milestones",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

# Timeline bar
timeline_y = Inches(2.0)
timeline_h = Inches(0.12)
bar_left = Inches(0.8)
bar_right = Inches(12.5)
bar_w = bar_right - bar_left

# Main timeline bar
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, timeline_y, bar_w, timeline_h)
shape.fill.solid()
shape.fill.fore_color.rgb = MED_GRAY
shape.line.fill.background()

# Milestones on timeline
milestones = [
    (0.04, "Feb 13", "MRD Review", ACCENT_BLUE),
    (0.20, "Mar 16", "Internal\nTesting", ACCENT_GREEN),
    (0.45, "Apr 20", "Demo\nCustomers", GOLD_ACCENT),
    (0.60, "May 25", "Product\nLaunch (GA)", ACCENT_BLUE),
    (0.78, "Jul 19", "Quality\nTargets Met", ACCENT_GREEN),
    (1.0, "Sep 16", "6 Paid\nPilots", RED_ACCENT),
]

for pct, date, label, color in milestones:
    x = bar_left + int(bar_w * pct) - Inches(0.15)

    # Dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, timeline_y - Inches(0.1), Inches(0.3), Inches(0.3)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    # Date above
    add_text_box(
        slide,
        x - Inches(0.3),
        timeline_y - Inches(0.55),
        Inches(0.9),
        Inches(0.4),
        date,
        font_size=12,
        bold=True,
        color=color,
        alignment=PP_ALIGN.CENTER,
    )

    # Label below
    add_text_box(
        slide,
        x - Inches(0.5),
        timeline_y + Inches(0.35),
        Inches(1.3),
        Inches(0.7),
        label,
        font_size=12,
        color=DARK_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

# Quality ramp arrow
add_text_box(
    slide,
    Inches(4.5),
    Inches(3.5),
    Inches(5),
    Inches(0.4),
    "90-day quality ramp (Apr 20 \u2192 Jul 19)",
    font_size=14,
    bold=True,
    color=ACCENT_GREEN,
    alignment=PP_ALIGN.CENTER,
)

# Feature readiness table
add_text_box(
    slide,
    Inches(0.8),
    Inches(4.2),
    Inches(11),
    Inches(0.4),
    "Feature Readiness",
    font_size=18,
    bold=True,
    color=DARK_TEXT,
)

ready_items = [
    ("Ready now: ", "Core RAG, access control, audit, React app, JWT auth"),
    ("By Mar 16 (internal): ", "AD/SSO integration"),
    (
        "By Apr 20 (demo): ",
        "Embeddable widget, usage reports, backup/restore, RAGAS evaluation, tiered XLSX",
    ),
]
add_bullet_list(
    slide,
    Inches(0.8),
    Inches(4.7),
    Inches(11),
    Inches(2),
    ready_items,
    font_size=15,
    color=DARK_TEXT,
)

# =============================================================================
# SLIDE 13: Success Metrics
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Success Metrics",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

# Business metrics
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5.5), Inches(4.0)
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(1.1),
    Inches(1.4),
    Inches(5),
    Inches(0.4),
    "Business Metrics (by Sep 16)",
    font_size=18,
    bold=True,
    color=ACCENT_BLUE,
)

biz_metrics = [
    ("6 ", "paid pilots closed"),
    ("$5,100 ", "MRR at cohort"),
    ("24 ", "qualified opportunities"),
    ("12 ", "proposals sent"),
    ("50% ", "close rate (proposal \u2192 pilot)"),
    ("\u2264$1,500 ", "CAC per pilot"),
    ("<$10,000 ", "total spend"),
]
add_bullet_list(
    slide,
    Inches(1.1),
    Inches(1.9),
    Inches(5),
    Inches(3),
    biz_metrics,
    font_size=16,
    color=DARK_TEXT,
)

# Product metrics
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.5), Inches(4.0)
)
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.fill.background()

add_text_box(
    slide,
    Inches(7.1),
    Inches(1.4),
    Inches(5),
    Inches(0.4),
    "Product Metrics (RAGAS)",
    font_size=18,
    bold=True,
    color=ACCENT_BLUE,
)

add_text_box(
    slide,
    Inches(7.1),
    Inches(1.9),
    Inches(2.2),
    Inches(0.3),
    "Launch:",
    font_size=14,
    bold=True,
    color=DARK_TEXT,
)
add_text_box(
    slide,
    Inches(9.3),
    Inches(1.9),
    Inches(2.5),
    Inches(0.3),
    "90-Day Target:",
    font_size=14,
    bold=True,
    color=DARK_TEXT,
)

quality_metrics = [
    ("Faithfulness: ", ">0.70 \u2192 >0.85"),
    ("Answer Relevancy: ", ">0.60 \u2192 >0.80"),
    ("Context Precision: ", ">0.50 \u2192 >0.70"),
    ("Context Recall: ", ">0.70 \u2192 >0.80"),
]
add_bullet_list(
    slide,
    Inches(7.1),
    Inches(2.3),
    Inches(5),
    Inches(1.8),
    quality_metrics,
    font_size=15,
    color=DARK_TEXT,
)

other_metrics = [
    ("System uptime: ", ">99.5%"),
    ("Response time (P95): ", "<10 seconds"),
]
add_bullet_list(
    slide,
    Inches(7.1),
    Inches(3.7),
    Inches(5),
    Inches(1.0),
    other_metrics,
    font_size=15,
    color=DARK_TEXT,
)

# Pilot KPIs at bottom
add_text_box(
    slide,
    Inches(0.8),
    Inches(5.5),
    Inches(11),
    Inches(0.4),
    "Per-Pilot KPIs",
    font_size=18,
    bold=True,
    color=DARK_TEXT,
)

pilot_kpis = [
    "Time-to-first-usable-answer  |  Reduction in document search time  |  Onboarding acceleration  |  Access-control correctness  |  Operator adoption"
]
add_bullet_list(
    slide,
    Inches(0.8),
    Inches(5.9),
    Inches(11),
    Inches(0.5),
    pilot_kpis,
    font_size=14,
    color=MED_GRAY,
)

# =============================================================================
# SLIDE 14: Risks
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0))

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Risks & Mitigations",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

risks = [
    (
        "HIGH",
        "No early proof points",
        "Enforce KPI instrumentation from day 1 of each pilot",
        RED_ACCENT,
    ),
    (
        "HIGH",
        "DIY seen as 'good enough'",
        "Sell time-to-value and support, not just technical capability",
        RED_ACCENT,
    ),
    (
        "MED",
        "Long security reviews",
        "Pre-package security brief and deployment checklist",
        GOLD_ACCENT,
    ),
    ("MED", "SMB budget friction", "Clear fixed $850/mo pricing. No surprise costs.", GOLD_ACCENT),
    ("MED", "DGX Spark supply constraints", "Hosted deployment as alternative path", GOLD_ACCENT),
    (
        "MED",
        "Model quality for domain",
        "90-day ramp with RAGAS evaluation. Swappable models.",
        GOLD_ACCENT,
    ),
]

y = Inches(1.2)
for level, risk, mitigation, color in risks:
    # Risk row
    row_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.85)
    )
    row_shape.fill.solid()
    row_shape.fill.fore_color.rgb = LIGHT_GRAY
    row_shape.line.fill.background()

    # Level badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y + Inches(0.18), Inches(0.9), Inches(0.45)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_text_box(
        slide,
        Inches(1.0),
        y + Inches(0.2),
        Inches(0.9),
        Inches(0.4),
        level,
        font_size=11,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Risk text
    txBox = slide.shapes.add_textbox(Inches(2.1), y + Inches(0.1), Inches(9.8), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_b = p.add_run()
    run_b.text = risk
    run_b.font.size = Pt(15)
    run_b.font.bold = True
    run_b.font.color.rgb = DARK_TEXT
    run_b.font.name = "Calibri"
    run_r = p.add_run()
    run_r.text = f"  \u2192  {mitigation}"
    run_r.font.size = Pt(14)
    run_r.font.color.rgb = MED_GRAY
    run_r.font.name = "Calibri"

    y += Inches(0.95)

# =============================================================================
# SLIDE 15: Decisions Needed
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide, top=Inches(0), color=RED_ACCENT)

add_text_box(
    slide,
    Inches(0.8),
    Inches(0.3),
    Inches(11),
    Inches(0.7),
    "Decisions Required",
    font_size=36,
    bold=True,
    color=DARK_TEXT,
)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.0),
    Inches(11),
    Inches(0.5),
    "6 decisions needed from this meeting to move forward:",
    font_size=20,
    color=DARK_TEXT,
)

decisions = [
    ("1. Product Name", '"AI Ready RAG" vs. alternatives — branding for market positioning'),
    ("2. Hardware Bundling", "Customer buys Spark separately, bundled pricing, or lease-to-own?"),
    ("3. Hosted Infrastructure", "Cloud provider, VM sizing, cost model per customer instance"),
    ("4. Post-Pilot Pricing", "Subscription tiers, support tiers, volume discounts"),
    ("5. Spanish Language", "Include in MVP, fast-follow, or evaluate only?"),
    ("6. Engineering Resources", "Team size, allocation, and sprint plan (Mar 16 \u2192 May 25)"),
]

y = Inches(1.7)
for title, desc in decisions:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    card.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.1), Inches(10.8), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_b = p.add_run()
    run_b.text = title
    run_b.font.size = Pt(17)
    run_b.font.bold = True
    run_b.font.color.rgb = RED_ACCENT
    run_b.font.name = "Calibri"
    run_r = p.add_run()
    run_r.text = f"  —  {desc}"
    run_r.font.size = Pt(15)
    run_r.font.color.rgb = DARK_TEXT
    run_r.font.name = "Calibri"

    y += Inches(0.9)

# Additional placeholders
add_text_box(
    slide,
    Inches(0.8),
    Inches(7.0),
    Inches(11),
    Inches(0.35),
    "Additional items to complete: TAM/SAM/SOM, financial model, customer quotes (post-Mar 16), partnership strategy",
    font_size=13,
    color=MED_GRAY,
    alignment=PP_ALIGN.CENTER,
)


# =============================================================================
# Save
# =============================================================================
output_path = "/home/jjob/projects/VE-RAG-System/docs/MRD-AI-Ready-RAG-v1.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
