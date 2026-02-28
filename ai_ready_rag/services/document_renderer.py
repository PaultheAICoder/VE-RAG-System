"""DocumentRenderer — core PDF and PPTX generation service.

Produces formatted output from structured data dicts.
CA automation services (board presentations, renewal packets, letters)
call this service rather than handling rendering themselves.

Dependencies (optional — graceful fallback if not installed):
  PDF: reportlab>=4.0.0
  PPTX: python-pptx>=0.6.23
"""

from __future__ import annotations

import logging
import os
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


@dataclass
class RenderResult:
    output_path: str
    format: str  # "pdf" | "pptx" | "html"
    file_size_bytes: int
    page_count: int | None = None
    slide_count: int | None = None


class DocumentRenderer:
    """Core renderer for structured document output.

    Supports PDF (via reportlab) and PPTX (via python-pptx).
    Falls back to plain HTML/text if libraries not installed.

    All output is written to temp files; caller is responsible for cleanup.
    """

    def __init__(self, output_dir: str | None = None) -> None:
        self._output_dir = output_dir or tempfile.gettempdir()

    def render_pdf(self, content: dict[str, Any], filename: str = "output.pdf") -> RenderResult:
        """Render structured content to PDF.

        content dict expected keys:
          - title: str
          - sections: list of {"heading": str, "body": str | list[str]}
          - metadata: dict (optional)
        """
        output_path = os.path.join(self._output_dir, filename)

        try:
            return self._render_pdf_reportlab(content, output_path)
        except ImportError:
            logger.warning("document_renderer.reportlab_not_installed", extra={"fallback": "html"})
            return self._render_html_fallback(content, output_path.replace(".pdf", ".html"))

    def _render_pdf_reportlab(self, content: dict[str, Any], output_path: str) -> RenderResult:
        """Render PDF using reportlab."""
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Title
        title = content.get("title", "Document")
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.2 * inch))

        # Sections
        for section in content.get("sections", []):
            heading = section.get("heading", "")
            if heading:
                story.append(Paragraph(heading, styles["Heading2"]))

            body = section.get("body", "")
            if isinstance(body, list):
                for item in body:
                    story.append(Paragraph(f"• {item}", styles["Normal"]))
            elif isinstance(body, str) and body:
                story.append(Paragraph(body, styles["Normal"]))

            story.append(Spacer(1, 0.1 * inch))

        doc.build(story)

        file_size = os.path.getsize(output_path)
        logger.info("document_renderer.pdf_created", extra={"path": output_path, "size": file_size})
        return RenderResult(
            output_path=output_path,
            format="pdf",
            file_size_bytes=file_size,
            page_count=None,  # reportlab doesn't easily expose page count
        )

    def render_pptx(self, content: dict[str, Any], filename: str = "output.pptx") -> RenderResult:
        """Render board presentation slides to PPTX.

        content dict expected from BoardPresentationPacket.to_dict():
          - account_name: str
          - meeting_date: str
          - slides: list of {"type": str, "title": str, "data": dict}
        """
        output_path = os.path.join(self._output_dir, filename)

        try:
            return self._render_pptx_python_pptx(content, output_path)
        except ImportError:
            logger.warning(
                "document_renderer.python_pptx_not_installed", extra={"fallback": "html"}
            )
            return self._render_html_fallback(content, output_path.replace(".pptx", ".html"))

    def _render_pptx_python_pptx(self, content: dict[str, Any], output_path: str) -> RenderResult:
        """Render PPTX using python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt  # noqa: F401

        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = content.get("account_name", "Insurance Review")
        if slide.placeholders[1]:
            slide.placeholders[1].text = content.get("meeting_date", "")

        # Content slides
        bullet_layout = prs.slide_layouts[1]
        for slide_data in content.get("slides", []):
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = slide_data.get("title", "")

            # Add data as bullet points
            data = slide_data.get("data", {})
            body_text = self._flatten_slide_data(data)
            if body_text and slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                for line in body_text[:10]:  # max 10 bullets per slide
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0

        prs.save(output_path)

        file_size = os.path.getsize(output_path)
        slide_count = len(prs.slides)
        logger.info(
            "document_renderer.pptx_created",
            extra={"path": output_path, "slides": slide_count},
        )
        return RenderResult(
            output_path=output_path,
            format="pptx",
            file_size_bytes=file_size,
            slide_count=slide_count,
        )

    def _flatten_slide_data(self, data: dict[str, Any]) -> list[str]:
        """Convert a data dict to a list of readable bullet strings."""
        bullets = []
        for key, value in data.items():
            if isinstance(value, list):
                for item in value:
                    if isinstance(item, dict):
                        bullets.append(
                            ", ".join(f"{k}: {v}" for k, v in item.items() if v is not None)
                        )
                    else:
                        bullets.append(str(item))
            elif isinstance(value, (int, float, str, bool)) and value is not None:
                label = key.replace("_", " ").title()
                if isinstance(value, float):
                    bullets.append(
                        f"{label}: ${value:,.2f}" if "usd" in key else f"{label}: {value:.2f}"
                    )
                else:
                    bullets.append(f"{label}: {value}")
        return bullets

    def _render_html_fallback(self, content: dict[str, Any], output_path: str) -> RenderResult:
        """Plain HTML fallback when PDF/PPTX libraries not installed."""
        title = content.get("title") or content.get("account_name", "Document")
        sections = content.get("sections") or content.get("slides", [])

        lines = [
            f"<html><head><title>{title}</title></head><body>",
            f"<h1>{title}</h1>",
        ]

        for section in sections:
            heading = section.get("heading") or section.get("title", "")
            if heading:
                lines.append(f"<h2>{heading}</h2>")
            body = section.get("body") or section.get("data", {})
            if isinstance(body, str):
                lines.append(f"<p>{body}</p>")
            elif isinstance(body, list):
                lines.append("<ul>")
                for item in body:
                    lines.append(f"<li>{item}</li>")
                lines.append("</ul>")
            elif isinstance(body, dict):
                lines.append("<ul>")
                for k, v in body.items():
                    lines.append(f"<li><strong>{k}</strong>: {v}</li>")
                lines.append("</ul>")

        lines.append("</body></html>")
        html = "\n".join(lines)

        Path(output_path).write_text(html, encoding="utf-8")
        file_size = os.path.getsize(output_path)

        return RenderResult(
            output_path=output_path,
            format="html",
            file_size_bytes=file_size,
        )
