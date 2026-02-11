"""Simple chunker for laptop development - lightweight dependencies."""

import logging
from pathlib import Path
from typing import Any

from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)


class SimpleChunker:
    """Lightweight chunker using RecursiveCharacterTextSplitter.

    Implements ChunkerProtocol for laptop profile.
    Uses minimal dependencies (no Docling, no Tesseract).
    """

    # Max pixels before resize to prevent OOM (20 megapixels)
    MAX_IMAGE_PIXELS = 20_000_000

    def __init__(
        self,
        chunk_size: int = 2000,
        chunk_overlap: int = 200,
        enable_ocr: bool = False,
        ocr_language: str = "eng",
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.enable_ocr = enable_ocr
        self.ocr_language = ocr_language
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len,
        )

    def chunk_document(
        self,
        file_path: str,
        metadata: dict[str, Any] | None = None,
    ) -> list[dict[str, Any]]:
        """Extract text from document and chunk it.

        Args:
            file_path: Path to document file
            metadata: Optional metadata to include in chunks

        Returns:
            List of chunk dicts with text, chunk_index, page_number, section
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        text = self._extract_text(path)
        chunks = self.chunk_text(text, source=path.name)

        # Add any provided metadata
        if metadata:
            for chunk in chunks:
                chunk.update(metadata)

        logger.info(f"Chunked {path.name}: {len(chunks)} chunks")
        return chunks

    def chunk_text(
        self,
        text: str,
        source: str = "unknown",
    ) -> list[dict[str, Any]]:
        """Chunk plain text.

        Args:
            text: Text content to chunk
            source: Source identifier for metadata

        Returns:
            List of chunk dicts
        """
        if not text or not text.strip():
            return []

        chunks = self.splitter.split_text(text)

        return [
            {
                "text": chunk,
                "chunk_index": i,
                "source": source,
                "page_number": None,  # Not available with simple extraction
                "section": None,
            }
            for i, chunk in enumerate(chunks)
        ]

    def _extract_text(self, path: Path) -> str:
        """Extract text from file based on extension.

        Supports: PDF, DOCX, TXT, MD, CSV, HTML, XLSX, PPTX,
                  PNG, JPG, JPEG, TIFF, TIF (OCR), EML, MSG (email)
        """
        ext = path.suffix.lower()

        if ext == ".pdf":
            return self._extract_pdf(path)
        elif ext == ".docx":
            return self._extract_docx(path)
        elif ext in (".txt", ".md"):
            return self._extract_text_file(path)
        elif ext == ".csv":
            return self._extract_csv(path)
        elif ext in (".html", ".htm"):
            return self._extract_html(path)
        elif ext == ".xlsx":
            return self._extract_xlsx(path)
        elif ext == ".pptx":
            return self._extract_pptx(path)
        elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif"):
            return self._extract_image(path)
        elif ext == ".eml":
            return self._extract_eml(path)
        elif ext == ".msg":
            return self._extract_msg(path)
        else:
            # Try as plain text
            try:
                return self._extract_text_file(path)
            except Exception as e:
                raise ValueError(f"Unsupported file type: {ext}") from e

    def _extract_text_file(self, path: Path) -> str:
        """Extract text from plain text file."""
        return path.read_text(encoding="utf-8", errors="replace")

    def _extract_pdf(self, path: Path) -> str:
        """Extract text from PDF using pypdf."""
        try:
            from pypdf import PdfReader
        except ImportError as e:
            raise ImportError("pypdf required for PDF extraction: pip install pypdf") from e

        reader = PdfReader(path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    def _extract_docx(self, path: Path) -> str:
        """Extract text from DOCX using python-docx.

        Extracts paragraphs, tables, headers, and footers.
        """
        try:
            from docx import Document
        except ImportError as e:
            raise ImportError(
                "python-docx required for DOCX extraction: pip install python-docx"
            ) from e

        doc = Document(path)
        content = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)

        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                if any(row_text):
                    table_text.append(" | ".join(row_text))
            if table_text:
                content.append("\n".join(table_text))

        # Extract headers and footers
        for section in doc.sections:
            if section.header and section.header.paragraphs:
                for para in section.header.paragraphs:
                    if para.text.strip():
                        content.append(para.text)
            if section.footer and section.footer.paragraphs:
                for para in section.footer.paragraphs:
                    if para.text.strip():
                        content.append(para.text)

        return "\n\n".join(content)

    def _extract_csv(self, path: Path) -> str:
        """Extract text from CSV."""
        import csv

        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = [", ".join(row) for row in reader]
        return "\n".join(rows)

    def _extract_html(self, path: Path) -> str:
        """Extract text from HTML."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            # Fallback: strip HTML tags with regex
            import re

            html = path.read_text(encoding="utf-8", errors="replace")
            text = re.sub(r"<[^>]+>", " ", html)
            return re.sub(r"\s+", " ", text).strip()

        html = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style"]):
            element.decompose()

        return soup.get_text(separator="\n", strip=True)

    def _extract_xlsx(self, path: Path) -> str:
        """Extract text from Excel file."""
        try:
            import openpyxl
        except ImportError as e:
            raise ImportError("openpyxl required for XLSX extraction: pip install openpyxl") from e

        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(cells):
                    rows.append(", ".join(cells))
            if rows:
                sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)

    def _extract_pptx(self, path: Path) -> str:
        """Extract text from PowerPoint file."""
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError(
                "python-pptx required for PPTX extraction: pip install python-pptx"
            ) from e

        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(slides)

    def _extract_image(self, path: Path) -> str:
        """Extract text from image using OCR.

        Requires enable_ocr=True. Raises ValueError if OCR is disabled.
        Large images (>20MP) are resized to prevent OOM.
        """
        if not self.enable_ocr:
            raise ValueError(
                f"Cannot process image file '{path.name}': OCR is disabled. "
                "Enable OCR (enable_ocr=True) to extract text from images."
            )

        try:
            import pytesseract
            from PIL import Image
        except ImportError as e:
            raise ImportError(
                "pytesseract and Pillow required for image extraction: "
                "pip install pytesseract Pillow"
            ) from e

        image = Image.open(path)

        # Resize large images to prevent OOM
        width, height = image.size
        pixels = width * height
        if pixels > self.MAX_IMAGE_PIXELS:
            scale = (self.MAX_IMAGE_PIXELS / pixels) ** 0.5
            new_width = int(width * scale)
            new_height = int(height * scale)
            logger.warning(
                "Image %s is %dx%d (%dMP), resizing to %dx%d",
                path.name,
                width,
                height,
                pixels // 1_000_000,
                new_width,
                new_height,
            )
            image = image.resize((new_width, new_height), Image.LANCZOS)

        text = pytesseract.image_to_string(image, lang=self.ocr_language)
        return text.strip()

    def _extract_eml(self, path: Path) -> str:
        """Extract text from .eml email files.

        Extracts headers (From, To, Date, Subject) and body text.
        Prefers plain text body; falls back to HTML with tag stripping.
        """
        import email
        from email import policy

        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)

        parts = []
        parts.append(f"From: {msg.get('From', 'Unknown')}")
        parts.append(f"To: {msg.get('To', 'Unknown')}")
        parts.append(f"Date: {msg.get('Date', 'Unknown')}")
        parts.append(f"Subject: {msg.get('Subject', 'No Subject')}")
        parts.append("")

        body = msg.get_body(preferencelist=("plain", "html"))
        if body:
            content = body.get_content()
            if body.get_content_type() == "text/html":
                content = self._strip_html_tags(content)
            parts.append(content)

        return "\n".join(parts)

    def _extract_msg(self, path: Path) -> str:
        """Extract text from Outlook .msg files.

        Extracts headers (From, To, Date, Subject) and body text.
        Binary attachments are skipped.
        """
        try:
            import extract_msg
        except ImportError as e:
            raise ImportError(
                "extract-msg required for .msg extraction: pip install extract-msg"
            ) from e

        msg = extract_msg.Message(path)
        try:
            parts = []
            parts.append(f"From: {msg.sender or 'Unknown'}")
            parts.append(f"To: {msg.to or 'Unknown'}")
            parts.append(f"Date: {msg.date or 'Unknown'}")
            parts.append(f"Subject: {msg.subject or 'No Subject'}")
            parts.append("")
            parts.append(msg.body or "")
            return "\n".join(parts)
        finally:
            msg.close()

    def _strip_html_tags(self, html: str) -> str:
        """Strip HTML tags from content, returning plain text."""
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html, "html.parser")
            for element in soup(["script", "style"]):
                element.decompose()
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            import re

            text = re.sub(r"<[^>]+>", " ", html)
            return re.sub(r"\s+", " ", text).strip()
